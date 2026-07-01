"""
report_engine.py
----------------
All report-generation logic for the Creative Audit PPTX pipeline.
No Colab or ipywidgets dependencies — imported by app.py (Streamlit).

API keys are read from environment variables or Streamlit secrets.
Set them in your Streamlit Cloud dashboard under Settings → Secrets:
    SMARTSHEET_API_KEY = "..."
    GEMINI_API_KEY     = "..."
"""

import os, sys, math, logging, gc
from pathlib import Path

import smartsheet
import requests
from PIL import Image
Image.MAX_IMAGE_PIXELS = None

import matplotlib
matplotlib.use("Agg")
matplotlib.rcParams["figure.max_open_warning"] = 5
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

import google.generativeai as genai

# ═══════════════════════════════════════════════════════════════
# API KEYS  —  read from env / Streamlit secrets (never hardcode)
# ═══════════════════════════════════════════════════════════════
def _get_secret(name: str) -> str:
    """Try Streamlit secrets first, fall back to env var."""
    try:
        import streamlit as st
        return st.secrets[name]
    except Exception:
        return os.getenv(name, "")

SMARTSHEET_API_KEY = _get_secret("SMARTSHEET_API_KEY")
GEMINI_API_KEY     = _get_secret("GEMINI_API_KEY")
GEMINI_MODEL       = "gemini-3.1-pro-preview"
SHEET_ID           = 84909979815812

# ═══════════════════════════════════════════════════════════════
# COLUMN NAMES
# ═══════════════════════════════════════════════════════════════
COL_CATEGORY = "Category"
COL_BRAND    = "Brand"
COL_ITEMTYPE = "Item Type"
COL_BG1UL    = "BG / 1UL"
COL_CLUSTER  = "BU / Cluster"
COL_COUNTRY  = "Country"
COL_YEAR     = "YEAR"
COL_MONTH    = "Month"
COL_SCORE    = "Score"
COL_REFCODE  = "Ref Code"
COL_BIZGROUP = "Business Group"

COL_CRITERIA = [
    "1. Clear Logo @3M?",
    "2. Brand Colours & DBA's",
    "3. Four Elements",
    "4. Good Contrast",
    "5. Face & eyes @3M?",
    '6. "NEW" @3M?',
    "7. Clear Visual Hierarchy",
    "8. Images other than product?",
    "9. Use Curves",
    "10. Vis L, TXT R",
    "11. Looking Prod / Message",
    "12. Avoid CAPITALS",
    "13. Buy Message",
    "14. 7 Words or fewer",
    "15. Instructional Verb?",
]

CRIT_SHORT_LABELS = [
    "Logo", "DBA", "4 Elements", "Contrast", "Face & Eye",
    "New?", "Hierarchy", "Images", "Curves", "Vis Left",
    "Gaze", "Capitals", "Buy Msg", "7 Words", "Imperative",
]

MONTH_ORDER = [
    "Jan", "Feb", "Mar", "Apr", "May", "Jun",
    "Jul", "Aug", "Sep", "Oct", "Nov", "Dec",
]
MONTH_INDEX = {m: i for i, m in enumerate(MONTH_ORDER)}

# ═══════════════════════════════════════════════════════════════
# PATHS
# ═══════════════════════════════════════════════════════════════
DOWNLOAD_DIR        = Path("/tmp/captivate_images")
TITLE_TEMPLATE_PATH = Path("title_template.pdf")   # place alongside app.py

# ═══════════════════════════════════════════════════════════════
# DESIGN TOKENS  —  Option C: white bg, purple accents, white header
# ═══════════════════════════════════════════════════════════════
C_BG         = RGBColor(0xFF, 0xFF, 0xFF)
C_HEADER     = RGBColor(0xFF, 0xFF, 0xFF)
C_HEADER_BDR = RGBColor(0xE0, 0xE0, 0xE0)
C_STRIPE     = RGBColor(0x3C, 0x2A, 0x8C)
C_CARD       = RGBColor(0xFF, 0xFF, 0xFF)
C_CARD_BG    = RGBColor(0xF7, 0xF5, 0xFF)
C_IMG_BG     = RGBColor(0xF2, 0xF2, 0xF2)
C_DIVIDER    = RGBColor(0xE8, 0xE8, 0xE8)
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK      = RGBColor(0x1A, 0x1A, 0x1A)
C_PURPLE     = RGBColor(0x3C, 0x2A, 0x8C)
C_LABEL      = RGBColor(0x88, 0x88, 0x88)
C_TEAL       = RGBColor(0x00, 0x99, 0x99)
C_RED_SCORE  = RGBColor(0xC0, 0x00, 0x00)
C_AMBER      = RGBColor(0xE0, 0x8C, 0x00)
C_GOLD       = C_TEAL
C_PINK       = C_RED_SCORE
C_GREEN      = C_TEAL
C_ORANGE_LBL = C_AMBER
C_RED_LBL    = C_RED_SCORE

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ═══════════════════════════════════════════════════════════════
# LOGGING
# ═══════════════════════════════════════════════════════════════
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s  %(levelname)-8s  %(message)s",
    datefmt="%H:%M:%S",
)
log = logging.getLogger(__name__)


# ───────────────────────────────────────────────────────────────
# SHAPE HELPERS
# ───────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line=None, lw=0, radius=0):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(lw)
    else:
        shape.line.fill.background()
    if radius:
        spPr     = shape._element.find(qn("p:spPr"))
        prstGeom = spPr.find(qn("a:prstGeom"))
        if prstGeom is not None:
            prstGeom.set("prst", "roundRect")
            avLst = prstGeom.find(qn("a:avLst"))
            if avLst is None:
                avLst = etree.SubElement(prstGeom, qn("a:avLst"))
            gd = etree.SubElement(avLst, qn("a:gd"))
            gd.set("name", "adj")
            gd.set("fmla", f"val {radius}")
    return shape


def add_text(
    slide, text, x, y, w, h,
    font="Poppins", size=12, bold=False, italic=False,
    color=None, align=PP_ALIGN.LEFT, anchor="t",
    wrap=True, spacing=0,
):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    bp = tf._txBody.find(qn("a:bodyPr"))
    if bp is None:
        bp = etree.SubElement(tf._txBody, qn("a:bodyPr"))
    bp.set("anchor", anchor)
    if spacing:
        bp.set("spc", str(spacing))
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name   = font
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txb


def score_to_label(s):
    if s >= 85:
        return "UNMISSABLE", C_TEAL
    if s >= 55:
        return "WALLPAPER", C_AMBER
    return "MISSABLE", C_RED_SCORE


# ───────────────────────────────────────────────────────────────
# SPEEDOMETER
# ───────────────────────────────────────────────────────────────

def generate_speedometer(score_str: str, out_path: Path) -> Path:
    try:
        s = score_str.strip().replace("%", "")
        score = float(s)
        if score <= 1.0:
            score *= 100
    except Exception:
        score = 0.0
    MAX_SCORE = 120.0
    score = max(0.0, min(MAX_SCORE, score))

    bg_hex = "#FFFFFF"
    fig, ax = plt.subplots(figsize=(3.5, 2.2), facecolor=bg_hex)
    ax.set_facecolor(bg_hex)
    ax.set_aspect("equal")
    ax.axis("off")

    n_steps = 120
    for i in range(n_steps):
        t      = i / n_steps
        theta1 = 180 - t * 180
        theta2 = 180 - (t + 1 / n_steps) * 180
        if t < 0.45:
            r, g, b = 0.95, t * 1.33, 0.1
        elif t < 0.70:
            tt = (t - 0.45) / 0.25
            r  = 0.95 - tt * 0.35
            g  = 0.60 + tt * 0.35
            b  = 0.1
        else:
            tt = (t - 0.70) / 0.30
            r  = 0.60 - tt * 0.55
            g  = min(1.0, 0.95 + tt * 0.05)
            b  = 0.1
        r, g, b = max(0, min(1, r)), max(0, min(1, g)), max(0, min(1, b))
        ax.add_patch(mpatches.Wedge(
            (0, 0), 1.0, theta2, theta1,
            width=0.28, facecolor=(r, g, b), edgecolor="none",
        ))

    ax.add_patch(plt.Circle((0, 0), 0.70, color=bg_hex, zorder=3))

    na = math.radians(180 - (score / MAX_SCORE) * 180)
    nx, ny = 0.62 * math.cos(na), 0.62 * math.sin(na)
    ax.plot([0, nx], [0, ny], color="#1A1A1A", lw=2.0, zorder=5, solid_capstyle="round")
    ax.add_patch(plt.Circle((0, 0), 0.05, color="#1A1A1A", zorder=6))
    ax.add_patch(plt.Circle((nx, ny), 0.035, color="#1A1A1A", zorder=6))

    lbl, lbl_color_rgb = score_to_label(score)
    hex_color = "#{:02X}{:02X}{:02X}".format(*lbl_color_rgb)
    ax.text(0, -0.45, f"{score:.0f}%", ha="center", va="center",
            fontsize=22, fontweight="bold", color=hex_color, zorder=7)
    ax.text(-1.08, -0.08, "MISSABLE",   ha="center", va="top",    fontsize=5.5, color="#888888")
    ax.text( 1.08, -0.08, "UNMISSABLE", ha="center", va="top",    fontsize=5.5, color="#888888")
    ax.text( 0,     1.08, "WALLPAPER",  ha="center", va="bottom", fontsize=5.5, color="#888888")

    ax.set_xlim(-1.25, 1.25)
    ax.set_ylim(-0.65, 1.2)
    plt.tight_layout(pad=0.1)
    plt.savefig(str(out_path), dpi=150, bbox_inches="tight", facecolor=bg_hex, format="png")
    plt.close(fig)
    return out_path


# ───────────────────────────────────────────────────────────────
# SMARTSHEET CONNECTION
# ───────────────────────────────────────────────────────────────

def connect_to_smartsheet():
    key = SMARTSHEET_API_KEY or os.getenv("SMARTSHEET_API_KEY")
    if not key:
        raise ValueError("SMARTSHEET_API_KEY not set. Add it to Streamlit secrets.")
    client = smartsheet.Smartsheet(key)
    client.errors_as_exceptions(True)
    log.info("Connected to Smartsheet.")
    return client


def load_sheet(client):
    log.info("Fetching sheet %s …", SHEET_ID)
    sheet   = client.Sheets.get_sheet(SHEET_ID)
    col_map = {col.title: col.id for col in sheet.columns}
    return sheet, col_map


# ───────────────────────────────────────────────────────────────
# DYNAMIC OPTION EXTRACTION
# ───────────────────────────────────────────────────────────────

def _cell_val(cell_map, col_map, col_name):
    cid  = col_map.get(col_name)
    cell = cell_map.get(cid) if cid else None
    return str(cell.value).strip() if (cell and cell.value is not None) else ""


def extract_unique_values(sheet, col_map):
    buckets = {
        COL_BIZGROUP: set(), COL_CATEGORY: set(), COL_BRAND: set(),
        COL_ITEMTYPE: set(), COL_BG1UL: set(), COL_CLUSTER: set(),
        COL_COUNTRY:  set(), COL_YEAR:  set(), COL_MONTH:  set(),
    }
    for row in sheet.rows:
        cm = {cell.column_id: cell for cell in row.cells}
        for col_name in buckets:
            v = _cell_val(cm, col_map, col_name)
            if v:
                if col_name == COL_YEAR:
                    try:
                        v = str(int(float(v)))
                    except ValueError:
                        pass
                elif col_name != COL_MONTH:
                    v = v.title()
                buckets[col_name].add(v)

    def month_sort_key(m):
        try:
            return int(m.strip().split(".")[0].strip())
        except (ValueError, IndexError):
            pass
        m_lower = m.lower()
        for idx, short in enumerate(MONTH_ORDER):
            if short.lower() in m_lower:
                return idx
        return 99

    result = {}
    for col_name, vals in buckets.items():
        if col_name == COL_MONTH:
            result[col_name] = sorted(vals, key=month_sort_key)
        elif col_name == COL_YEAR:
            result[col_name] = sorted(vals, key=lambda y: int(y) if y.isdigit() else 0)
        else:
            result[col_name] = sorted(vals, key=str.upper)
    return result


# ───────────────────────────────────────────────────────────────
# HEADER TEXT BUILDER
# ───────────────────────────────────────────────────────────────

def build_header_text(filters: dict) -> str:
    parts = []
    hierarchy = [
        ("bizgroup", "Business Group"), ("category", "Category"),
        ("brand", "Brand"),            ("itemtype", "Item Type"),
        ("bg1ul", "BG / 1UL"),         ("cluster",  "BU / Cluster"),
        ("country", "Country"),         ("year",     "Year"),
        ("month", "Month"),
    ]
    for key, _ in hierarchy:
        val = filters.get(key)
        if val:
            parts.append(" + ".join(v.upper() for v in val) if isinstance(val, list) else str(val).upper())

    s_min = filters.get("score_min")
    s_max = filters.get("score_max")
    if s_min is not None and s_max is not None:
        parts.append(f"SCORE {s_min}-{s_max}%")
    elif s_min is not None:
        parts.append(f"SCORE {s_min}%+")
    elif s_max is not None:
        parts.append(f"SCORE UP TO {s_max}%")

    return "  |  ".join(parts) if parts else "ALL DATA"


# ───────────────────────────────────────────────────────────────
# FILTERED ROW RETRIEVAL
# ───────────────────────────────────────────────────────────────

def get_filtered_rows(sheet, col_map, filters):
    simple = {
        COL_CATEGORY: filters.get("category"),
        COL_BIZGROUP: filters.get("bizgroup"),
        COL_BRAND:    filters.get("brand"),
        COL_ITEMTYPE: filters.get("itemtype"),
        COL_BG1UL:    filters.get("bg1ul"),
        COL_CLUSTER:  filters.get("cluster"),
        COL_COUNTRY:  filters.get("country"),
    }
    year_val = filters.get("year")
    if year_val:
        simple[COL_YEAR] = (
            [str(float(int(y))) for y in year_val]
            if isinstance(year_val, list)
            else str(float(int(year_val)))
        )

    match_map = {}
    for col_name, val in simple.items():
        if val and col_name in col_map:
            match_map[col_map[col_name]] = val

    m_start = filters.get("month_start")
    m_end   = filters.get("month_end")
    month_range_active = bool(m_start and m_end)

    def _month_num(m_str):
        m_str = m_str.strip()
        try:
            return int(m_str.split(".")[0].strip()) - 1
        except (ValueError, IndexError):
            pass
        m_lower = m_str.lower()
        for idx, short in enumerate(MONTH_ORDER):
            if short.lower() in m_lower:
                return idx
        return 0

    valid_month_nums = set()
    if month_range_active:
        si = _month_num(m_start)
        ei = _month_num(m_end)
        valid_month_nums = set(range(min(si, ei), max(si, ei) + 1))
    month_col_id = col_map.get(COL_MONTH)

    results = []
    for row in sheet.rows:
        cm = {cell.column_id: cell for cell in row.cells}

        skip = False
        for cid, expected in match_map.items():
            raw_actual   = str(cm.get(cid) and cm[cid].value or "").strip()
            actual_title = raw_actual.title()
            if isinstance(expected, list):
                allowed = {str(e).title() for e in expected} | {str(e) for e in expected}
                if actual_title not in allowed and raw_actual not in allowed:
                    skip = True; break
            else:
                if actual_title != str(expected).title() and raw_actual != str(expected):
                    skip = True; break
        if skip:
            continue

        if month_range_active and month_col_id:
            month_val = str(cm.get(month_col_id) and cm[month_col_id].value or "").strip()
            if _month_num(month_val) not in valid_month_nums:
                continue

        def cval(name):
            return _cell_val(cm, col_map, name)

        raw = cval(COL_SCORE)
        try:
            score_float = float(raw) * 100
            score_pct   = f"{score_float:.0f}%"
        except Exception:
            score_float = 0
            score_pct   = raw

        score_min = filters.get("score_min")
        score_max = filters.get("score_max")
        if score_min is not None and score_float < score_min:
            continue
        if score_max is not None and score_float > score_max:
            continue

        rd = {
            "row_id":   row.id,
            "bizgroup": cval(COL_BIZGROUP),
            "brand":    cval(COL_BRAND),
            "itemtype": cval(COL_ITEMTYPE),
            "bg1ul":    cval(COL_BG1UL),
            "country":  cval(COL_COUNTRY),
            "category": cval(COL_CATEGORY),
            "cluster":  cval(COL_CLUSTER),
            "score":    score_pct,
            "ref_code": cval(COL_REFCODE),
        }
        for i, cn in enumerate(COL_CRITERIA, 1):
            rd[f"crit_{i}"] = cval(cn)
        results.append(rd)

    log.info("Found %d matching row(s).", len(results))
    return results


# ───────────────────────────────────────────────────────────────
# ATTACHMENTS
# ───────────────────────────────────────────────────────────────

def fetch_images_for_row(client, row_id):
    atts = client.Attachments.list_row_attachments(SHEET_ID, row_id)
    if not atts.data:
        return []
    DOWNLOAD_DIR.mkdir(parents=True, exist_ok=True)
    paths = []
    IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp", ".tiff"}
    image_atts = [a for a in atts.data if Path(a.name or "").suffix.lower() in IMAGE_EXTS]
    for att in image_atts:
        try:
            full  = client.Attachments.get_attachment(SHEET_ID, att.id)
            fname = att.name or f"row_{row_id}_{att.id}.jpg"
            dest  = DOWNLOAD_DIR / fname
            with requests.get(full.url, timeout=30, stream=True) as resp:
                resp.raise_for_status()
                dest.write_bytes(resp.content)
            with Image.open(dest) as img:
                img = img.convert("RGB")
                img.thumbnail((800, 600), Image.LANCZOS)
                cp = dest.with_stem(dest.stem + "_lb")
                img.save(cp, format="JPEG", quality=75, optimize=True)
            dest.unlink(missing_ok=True)
            paths.append(cp)
        except Exception as e:
            log.warning("Attachment error row %s: %s", row_id, e)
    return paths


# ───────────────────────────────────────────────────────────────
# SHARED SLIDE CHROME
# ───────────────────────────────────────────────────────────────

_HDR_H = Inches(0.62)
_M_L   = Inches(0.22)


def _ensure_logo():
    logo_path = DOWNLOAD_DIR / "logo_icon.png"
    if logo_path.exists():
        return logo_path
    DOWNLOAD_DIR.mkdir(parents=True, exist_ok=True)
    fig2, ax2 = plt.subplots(figsize=(1, 1), facecolor="none")
    ax2.set_aspect("equal"); ax2.axis("off")
    ax2.set_xlim(-1.1, 1.1); ax2.set_ylim(-1.1, 1.1)
    ax2.add_patch(mpatches.Wedge((0,0), 1.0,  90, 270, width=0.32, facecolor="#3C2A8C", edgecolor="none"))
    ax2.add_patch(mpatches.Wedge((0,0), 1.0, -90,  90, width=0.32, facecolor="#E08C00", edgecolor="none"))
    ax2.add_patch(plt.Circle((0,0), 0.66, color="white", zorder=3))
    ax2.add_patch(plt.Circle((0,0), 0.34, color="#3C2A8C", zorder=4))
    plt.tight_layout(pad=0)
    plt.savefig(str(logo_path), dpi=150, bbox_inches="tight", transparent=True, format="png")
    plt.close(fig2)
    return logo_path


def _draw_slide_chrome(slide, header_text):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C_BG)
    add_rect(slide, 0, 0, SLIDE_W, _HDR_H, fill=C_HEADER)
    add_rect(slide, 0, _HDR_H - Pt(1), SLIDE_W, Pt(1), fill=C_HEADER_BDR)
    logo_path = _ensure_logo()
    LOGO_S = Inches(0.44)
    logo_x = _M_L
    logo_y = (_HDR_H - LOGO_S) / 2
    slide.shapes.add_picture(str(logo_path), logo_x, logo_y, LOGO_S, LOGO_S)
    add_text(
        slide, header_text,
        logo_x + LOGO_S + Inches(0.14), 0,
        SLIDE_W - logo_x - LOGO_S - Inches(0.3), _HDR_H,
        font="Poppins", size=13, bold=True,
        color=C_PURPLE, align=PP_ALIGN.LEFT, anchor="ctr",
    )
    return _HDR_H + Inches(0.2)


# ───────────────────────────────────────────────────────────────
# CREATIVE SLIDE
# ───────────────────────────────────────────────────────────────

def build_slide(prs, row, image_paths, header_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C_BG)

    M   = Inches(0.22)
    M_L = Inches(0.22)

    HDR_H = Inches(0.62)
    add_rect(slide, 0, 0, SLIDE_W, HDR_H, fill=C_HEADER)
    add_rect(slide, 0, HDR_H - Pt(1), SLIDE_W, Pt(1), fill=C_HEADER_BDR)

    LOGO_S    = Inches(0.44)
    logo_x    = M_L
    logo_y    = (HDR_H - LOGO_S) / 2
    logo_path = _ensure_logo()
    slide.shapes.add_picture(str(logo_path), logo_x, logo_y, LOGO_S, LOGO_S)

    add_text(
        slide, header_text,
        logo_x + LOGO_S + Inches(0.14), 0,
        SLIDE_W - logo_x - LOGO_S - Inches(0.3), HDR_H,
        font="Poppins", size=13, bold=True,
        color=C_PURPLE, align=PP_ALIGN.LEFT, anchor="ctr",
    )

    CONTENT_TOP = HDR_H + M
    FOOTER_H    = Inches(1.82)
    FOOTER_TOP  = SLIDE_H - FOOTER_H - Inches(0.12)
    RIGHT_W     = Inches(3.9)
    IMG_W       = SLIDE_W - RIGHT_W - M_L - M * 2
    IMG_H       = FOOTER_TOP - CONTENT_TOP - M
    img_x       = M_L
    img_y       = CONTENT_TOP
    panel_x     = img_x + IMG_W + M
    panel_w     = RIGHT_W - M
    panel_h     = IMG_H

    add_rect(slide, img_x, img_y, IMG_W, IMG_H, fill=C_IMG_BG)

    if not image_paths:
        add_text(slide, "CREATIVE PREVIEW",
                 img_x, img_y, IMG_W, IMG_H,
                 font="Poppins", size=16, italic=True, color=C_LABEL,
                 align=PP_ALIGN.CENTER, anchor="ctr")
    else:
        for img_path in image_paths:
            try:
                with Image.open(img_path) as im:
                    iw, ih = im.size
                img_ratio = iw / ih
                box_ratio = IMG_W / IMG_H
                if img_ratio > box_ratio:
                    pic_w = IMG_W
                    pic_h = int(IMG_W / img_ratio)
                else:
                    pic_h = IMG_H
                    pic_w = int(IMG_H * img_ratio)
                pic_x = img_x + (IMG_W - pic_w) // 2
                pic_y = img_y + (IMG_H - pic_h) // 2
                slide.shapes.add_picture(str(img_path), pic_x, pic_y, pic_w, pic_h)
            except Exception as e:
                log.warning("Image insert failed: %s", e)
        if len(image_paths) > 1:
            add_text(slide,
                     f"[WARNING] {len(image_paths)} images stacked - delete unwanted in PowerPoint",
                     img_x + Inches(0.1), img_y + Inches(0.1),
                     Inches(4), Inches(0.28),
                     size=7.5, bold=True, color=C_AMBER)

    add_rect(slide, panel_x, img_y, panel_w, panel_h, fill=C_WHITE,
             line=C_HEADER_BDR, lw=0.5)

    px = panel_x + Inches(0.22)
    pw = panel_w  - Inches(0.44)

    info_h  = panel_h * 0.44
    item_h  = info_h / 3
    start_y = img_y + Inches(0.18)

    info_items = [
        ("BRAND",    row["brand"],    True,  15),
        ("COUNTRY",  row["country"],  True,  14),
        ("REF CODE", row["ref_code"], False, 10),
    ]
    for idx, (lbl, val, bold, vsize) in enumerate(info_items):
        iy = start_y + idx * item_h
        if idx > 0:
            add_rect(slide, px, iy - Inches(0.05), pw, Pt(0.5), fill=C_DIVIDER)
        add_text(slide, lbl, px, iy, pw, Inches(0.22),
                 font="Poppins", size=7.5, color=C_LABEL, anchor="t")
        add_text(slide, val, px, iy + Inches(0.22), pw, Inches(0.52),
                 font="Poppins", size=vsize, bold=bold,
                 color=C_PURPLE if bold else C_BLACK, anchor="t")

    score_y = img_y + info_h + Inches(0.12)
    add_rect(slide, px, score_y - Inches(0.04), pw, Pt(0.5), fill=C_DIVIDER)
    add_text(slide, "SCORE", px, score_y, pw, Inches(0.22),
             font="Poppins", size=7.5, color=C_LABEL, anchor="t", align=PP_ALIGN.CENTER)

    try:
        sp_path = DOWNLOAD_DIR / f"speedo_{row['row_id']}.png"
        generate_speedometer(row["score"], sp_path)
        available_h = (img_y + panel_h) - (score_y + Inches(0.26)) - Inches(0.45)
        speedo_ratio = 3.5 / 2.2
        sp_w = pw
        sp_h = sp_w / speedo_ratio
        if sp_h > available_h:
            sp_h = available_h
            sp_w = sp_h * speedo_ratio
        sp_x = panel_x + (panel_w - sp_w) / 2
        sp_y = score_y + Inches(0.26)
        slide.shapes.add_picture(str(sp_path), sp_x, sp_y, sp_w, sp_h)
        try:
            sp_path.unlink(missing_ok=True)
        except Exception:
            pass
        try:
            score_num = float(row["score"].replace("%", ""))
        except Exception:
            score_num = 0
        lbl_text, lbl_color = score_to_label(score_num)
        badge_w = Inches(1.4)
        badge_h = Inches(0.28)
        badge_x = panel_x + (panel_w - badge_w) / 2
        badge_y = sp_y + sp_h + Inches(0.04)
        add_rect(slide, badge_x, badge_y, badge_w, badge_h, fill=lbl_color, radius=50000)
        add_text(slide, lbl_text, badge_x, badge_y, badge_w, badge_h,
                 font="Poppins", size=9, bold=True, color=C_WHITE,
                 align=PP_ALIGN.CENTER, anchor="ctr")
    except Exception as e:
        log.warning("Speedometer error: %s", e)

    add_rect(slide, 0, FOOTER_TOP, SLIDE_W, FOOTER_H, fill=C_CARD_BG, line=C_HEADER_BDR, lw=0.3)
    add_rect(slide, 0, FOOTER_TOP, SLIDE_W, Pt(2), fill=C_STRIPE)

    ft_y = FOOTER_TOP + Inches(0.1)
    add_text(slide, "EVALUATION CRITERIA",
             M_L, ft_y, Inches(3), Inches(0.22),
             font="Poppins", size=7.5, bold=True, color=C_PURPLE, anchor="t")

    n        = 15
    ft_w     = SLIDE_W - M_L - M
    cell_w   = ft_w / n
    btn_s    = Inches(0.54)
    btn_y    = ft_y + Inches(0.28)
    lbl_bt_y = btn_y + btn_s + Inches(0.05)

    for i in range(n):
        cx   = M_L + cell_w * i
        val  = row.get(f"crit_{i+1}", "").strip().upper()
        is_y = val == "Y"
        btn_x = cx + (cell_w - btn_s) / 2
        add_rect(slide, btn_x, btn_y, btn_s, btn_s,
                 fill=C_TEAL if is_y else C_RED_SCORE, radius=50000)
        add_text(slide, val if val in ("Y", "N") else "-",
                 btn_x, btn_y, btn_s, btn_s,
                 font="Poppins", size=15, bold=True, color=C_WHITE,
                 align=PP_ALIGN.CENTER, anchor="ctr")
        add_text(slide, CRIT_SHORT_LABELS[i],
                 cx, lbl_bt_y, cell_w, Inches(0.26),
                 size=7, color=C_LABEL,
                 align=PP_ALIGN.CENTER, anchor="t", wrap=True)

    log.info("Slide built: %s / %s / %s", row["brand"], row["ref_code"], row["score"])


# ───────────────────────────────────────────────────────────────
# SUMMARY SLIDE
# ───────────────────────────────────────────────────────────────

def build_summary_slide(prs, rows, header_text, analysis_mode='average'):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    M = Inches(0.22)
    CONTENT_TOP = _draw_slide_chrome(slide, header_text + "  |  SUMMARY")
    total = len(rows)

    fail_counts = [0] * 15
    for row in rows:
        for i in range(15):
            if row.get(f"crit_{i+1}", "").strip().upper() == "N":
                fail_counts[i] += 1

    pass_counts_cards = [total - fc for fc in fail_counts]
    scores_for_card = []
    for _r in rows:
        try:
            scores_for_card.append(float(_r["score"].replace("%", "")))
        except Exception:
            pass
    avg_score_val = sum(scores_for_card) / len(scores_for_card) if scores_for_card else 0
    missable_n   = sum(1 for s in scores_for_card if s < 55)
    wallpaper_n  = sum(1 for s in scores_for_card if 55 <= s < 85)
    unmissable_n = sum(1 for s in scores_for_card if s >= 85)

    if analysis_mode == "volume":
        card_labels = [
            ("CREATIVES REVIEWED", str(total)),
            ("UNMISSABLE (85%+)",   str(unmissable_n)),
            ("WALLPAPER (55-84%)",  str(wallpaper_n)),
            ("MISSABLE (<55%)",     str(missable_n)),
        ]
    else:
        card_labels = [
            ("CREATIVES REVIEWED",     str(total)),
            ("MOST PASSED CRITERION",  CRIT_SHORT_LABELS[pass_counts_cards.index(max(pass_counts_cards))] if total else "-"),
            ("LEAST PASSED CRITERION", CRIT_SHORT_LABELS[pass_counts_cards.index(min(pass_counts_cards))] if total else "-"),
            ("AVG CREATIVE SCORE",     f"{avg_score_val:.0f}%" if scores_for_card else "-"),
        ]

    extra_cards = []
    brands_in_data    = list({r["brand"]   for r in rows if r.get("brand")})
    countries_in_data = list({r["country"] for r in rows if r.get("country")})

    for field, field_label_top, field_label_low in [
        ("brand",   "TOP BRAND",    "LOWEST BRAND"),
        ("country", "TOP COUNTRY",  "LOWEST COUNTRY"),
    ]:
        in_data = brands_in_data if field == "brand" else countries_in_data
        if len(in_data) > 1:
            fscores = {}
            for r in rows:
                v = r.get(field, "")
                if v:
                    try:
                        fscores.setdefault(v, []).append(float(r["score"].replace("%", "")))
                    except Exception:
                        pass
            if fscores:
                if analysis_mode == "volume":
                    top_v = max(fscores, key=lambda k: len(fscores[k]))
                    low_v = min(fscores, key=lambda k: len(fscores[k]))
                    extra_cards.append(("MOST CREATIVES",   f"{top_v.title()} ({len(fscores[top_v])})"))
                    extra_cards.append(("FEWEST CREATIVES",  f"{low_v.title()} ({len(fscores[low_v])})"))
                else:
                    top_v = max(fscores, key=lambda k: sum(fscores[k])/len(fscores[k]))
                    low_v = min(fscores, key=lambda k: sum(fscores[k])/len(fscores[k]))
                    extra_cards.append((field_label_top, f"{top_v.title()} ({sum(fscores[top_v])/len(fscores[top_v]):.0f}%)"))
                    extra_cards.append((field_label_low, f"{low_v.title()} ({sum(fscores[low_v])/len(fscores[low_v]):.0f}%)"))

    all_cards = card_labels + extra_cards[:4]
    n_cards = len(all_cards)
    rows_of_cards = [all_cards] if n_cards <= 4 else [all_cards[:(n_cards+1)//2], all_cards[(n_cards+1)//2:]]

    CARD_ROW_H   = Inches(0.82)
    CARD_ROW_GAP = Inches(0.08)
    for row_idx, card_row in enumerate(rows_of_cards):
        card_w  = (SLIDE_W - _M_L - M) / len(card_row)
        row_top = CONTENT_TOP + row_idx * (CARD_ROW_H + CARD_ROW_GAP)
        for idx, (lbl, val) in enumerate(card_row):
            cx = _M_L + card_w * idx
            add_rect(slide, cx + Inches(0.06), row_top,
                     card_w - Inches(0.12), CARD_ROW_H,
                     fill=C_WHITE, line=C_HEADER_BDR, lw=0.5)
            add_rect(slide, cx + Inches(0.06), row_top,
                     Inches(0.05), CARD_ROW_H, fill=C_STRIPE)
            val_size = 15 if len(val) > 12 else 19
            add_text(slide, lbl,
                     cx + Inches(0.22), row_top + Inches(0.1),
                     card_w - Inches(0.32), Inches(0.26),
                     font="Poppins", size=7.5, color=C_LABEL, align=PP_ALIGN.LEFT, anchor="t")
            add_text(slide, val,
                     cx + Inches(0.22), row_top + Inches(0.36),
                     card_w - Inches(0.32), Inches(0.38),
                     font="Poppins", size=val_size, bold=True, color=C_PURPLE,
                     align=PP_ALIGN.LEFT, anchor="t")

    CARD_H    = CARD_ROW_H * len(rows_of_cards) + CARD_ROW_GAP * (len(rows_of_cards) - 1)
    CHART_TOP = CONTENT_TOP + CARD_H + Inches(0.22)
    CHART_H   = SLIDE_H - CHART_TOP - Inches(0.2)
    CHART_W   = SLIDE_W - _M_L - M

    chart_path = DOWNLOAD_DIR / "summary_chart.png"
    bg_hex = "#FFFFFF"
    fig, ax = plt.subplots(figsize=(12, 4.2), facecolor=bg_hex)
    ax.set_facecolor(bg_hex)

    x           = np.arange(15)
    pass_counts = [total - fc for fc in fail_counts]
    pcts        = [pc / total * 100 if total else 0 for pc in pass_counts]
    fail_pcts   = [fc / total * 100 if total else 0 for fc in fail_counts]

    bars_pass = ax.bar(x, pcts,      color="#009999", width=0.6, zorder=3)
    bars_fail = ax.bar(x, fail_pcts, color="#C00000", width=0.6, zorder=3, bottom=pcts)

    for bp, bf, pct, fail_p, pc, fc in zip(bars_pass, bars_fail, pcts, fail_pcts, pass_counts, fail_counts):
        if pct >= 10:
            ax.text(bp.get_x() + bp.get_width()/2, pct/2,
                    f"{pct:.0f}%\n({pc})\nYes",
                    ha="center", va="center", fontsize=7, color="white", fontweight="bold", zorder=5)
        if fail_p >= 10:
            ax.text(bf.get_x() + bf.get_width()/2, pct + fail_p/2,
                    f"{fail_p:.0f}%\n({fc})\nNo",
                    ha="center", va="center", fontsize=7, color="white", fontweight="bold", zorder=5)

    ax.set_xticks(x)
    ax.set_xticklabels(CRIT_SHORT_LABELS, fontsize=7.5, color="#444444", rotation=30, ha="right")
    ax.set_yticks(range(0, 101, 20))
    ax.set_yticklabels([f"{v}%" for v in range(0, 101, 20)], fontsize=7.5, color="#444444")
    ax.set_ylim(0, 100.1)
    ax.set_ylabel("Pass Rate", color="#444444", fontsize=8)
    ax.tick_params(axis="both", colors="#444444", length=0)
    for spine in ax.spines.values():
        spine.set_visible(False)
    ax.yaxis.grid(True, color="#E8E8E8", linewidth=0.5, zorder=0)
    ax.set_axisbelow(True)
    plt.tight_layout(pad=0.4)
    plt.savefig(str(chart_path), dpi=150, bbox_inches="tight", facecolor=bg_hex, format="png")
    plt.close(fig)

    slide.shapes.add_picture(str(chart_path), _M_L, CHART_TOP, CHART_W, CHART_H)
    try:
        chart_path.unlink(missing_ok=True)
    except Exception:
        pass
    log.info("Summary slide built (%d rows).", total)


# ───────────────────────────────────────────────────────────────
# TITLE SLIDE
# ───────────────────────────────────────────────────────────────

def build_title_slide(prs, header_text):
    from datetime import datetime as _dt
    date_str = _dt.now().strftime("%B %Y")

    candidates = [
        Path("title_template.pdf"),
        Path("/tmp/title_template.pdf"),
    ]
    pdf_path = next((p for p in candidates if p.exists()), None)

    if not pdf_path:
        log.warning("title_template.pdf not found — skipping title slide.")
        return

    DOWNLOAD_DIR.mkdir(parents=True, exist_ok=True)
    img_out   = DOWNLOAD_DIR / "title_slide_bg.png"
    converted = False

    try:
        import fitz
        doc  = fitz.open(str(pdf_path))
        page = doc.load_page(0)
        pix  = page.get_pixmap(matrix=fitz.Matrix(2.78, 2.78))
        pix.save(str(img_out))
        doc.close()
        if img_out.exists():
            converted = True
    except Exception as e:
        log.warning("PyMuPDF failed: %s", e)

    if not converted:
        log.warning("Could not convert title PDF — skipping title slide.")
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=RGBColor(0x3D, 0x1A, 0x8C))
    slide.shapes.add_picture(str(img_out), 0, 0, SLIDE_W, SLIDE_H)

    add_text(slide, header_text,
             Inches(0.4), SLIDE_H * 0.635,
             SLIDE_W - Inches(0.8), Inches(0.75),
             font="Poppins", size=26, bold=True,
             color=C_WHITE, align=PP_ALIGN.LEFT, anchor="t")
    add_text(slide, date_str,
             Inches(0.4), SLIDE_H * 0.635 + Inches(0.82),
             Inches(5), Inches(0.4),
             font="Poppins", size=13,
             color=C_AMBER, align=PP_ALIGN.LEFT, anchor="t")
    try:
        img_out.unlink(missing_ok=True)
    except Exception:
        pass


# ───────────────────────────────────────────────────────────────
# GEMINI AI SYNOPSIS
# ───────────────────────────────────────────────────────────────

def generate_gemini_synopsis(rows, filters, fail_counts):
    try:
        genai.configure(api_key=GEMINI_API_KEY)
        model = genai.GenerativeModel(GEMINI_MODEL)
    except Exception as e:
        log.warning("Gemini init failed: %s", e)
        return None

    total       = len(rows)
    pass_counts = [total - fc for fc in fail_counts]

    criteria_lines = []
    for label, pc, fc in zip(CRIT_SHORT_LABELS, pass_counts, fail_counts):
        pct = pc / total * 100 if total else 0
        criteria_lines.append(f"  - {label}: {pct:.0f}% pass rate ({pc}/{total})")

    FILTER_DISPLAY = {
        "bizgroup": "Business Group", "category": "Category", "brand": "Brand",
        "itemtype": "Item Type",      "bg1ul":    "BG / 1UL", "cluster": "BU / Cluster",
        "country":  "Country",        "year":     "Year",     "month":   "Month",
    }
    context_parts = []
    for k, v in filters.items():
        if k not in ("score_min", "score_max", "month_start", "month_end", "analysis_mode"):
            label = FILTER_DISPLAY.get(k, k.title())
            context_parts.append(f"{label}: {', '.join(v) if isinstance(v, list) else v}")
    context_str = ", ".join(context_parts) if context_parts else "All Data"

    scores    = []
    for row in rows:
        try:
            scores.append(float(row["score"].replace("%", "")))
        except Exception:
            pass
    avg_score  = sum(scores) / len(scores) if scores else 0
    top_scores = sorted(scores, reverse=True)[:3]
    low_scores = sorted(scores)[:3]

    prompt = f"""You are a world-class creative effectiveness analyst specialising in PHYSICAL RETAIL and IN-STORE shopper marketing for UNILEVER brands. Your expertise spans behavioural science, shopper psychology, and the most current peer-reviewed marketing effectiveness research including Kahneman (System 1/2 processing), Binet & Field (effectiveness), Byron Sharp and the Ehrenberg-Bass Institute (mental and physical availability, distinctiveness), and leading shopper marketing science.

UNILEVER BRAND CONTEXT:
- The creatives belong to well-known FMCG brands competing in physical retail (e.g. Dove, Axe, Lynx, Hellmann's, Persil, Sure, TRESemme, Simple, Comfort, Knorr, Wall's etc.)
- Never refer to "Unilever" by name - always refer to the specific brand name(s) in the data instead
- These are large, established FMCG brands where shelf standout and instant brand recognition are commercially critical

FIXED ASSUMPTIONS:
- DISTANCE: The shopper is always viewing from 3 metres away
- ATTENTION WINDOW: The shopper gives 2 seconds of attention maximum before moving on
- COGNITIVE CAPACITY: The human brain can only process 4 visual elements within that 2-second window
- COPY LIMIT: The shopper can only read 7 words in that 2-second window
- MEDIA TYPE: Always static media (print, POS, shelf materials)
- ENVIRONMENT: Always in-store physical retail
- Never use: "digital shelf", "scrolling", "feed", "online", "thumb-stopping", "swipe", "click", "screen"
- Always use: "shelf", "fixture", "in-store", "retail environment", "shopper", "aisle", "point of purchase", "3-metre distance", "2-second window"

SCORING FRAMEWORK:
- Above 84% = UNMISSABLE
- 55% to 84% = WALLPAPER
- Below 55% = MISSABLE

Banned words: "urgent", "urgently", "ruthlessly", "critical failure", "devastating", "alarming", "catastrophic", "brutal", "slash", "destroy", "eliminate"

CONTEXT: {context_str}
TOTAL CREATIVES REVIEWED: {total}
AVERAGE CREATIVE SCORE: {avg_score:.0f}%
TOP 3 SCORES: {', '.join(f'{s:.0f}%' for s in top_scores)}
LOWEST 3 SCORES: {', '.join(f'{s:.0f}%' for s in low_scores)}

CRITERIA PASS RATES (15 criteria):
{chr(10).join(criteria_lines)}

Write exactly 4 sections with these exact headings:

EXECUTIVE SUMMARY
2-3 sentences. Honest, calibrated assessment. Reference the average score and tier.

KEY STRENGTHS
Exactly 3 bullet points starting with •
Label in CAPS: then 1-2 sentences grounded in behavioural science. Reference pass rate %.

AREAS FOR IMPROVEMENT
Exactly 3 bullet points starting with •
Label in CAPS: then 1-2 sentences on why it matters behaviourally. Reference pass rate %.

PRIORITY RECOMMENDATIONS
Exactly 3 bullet points starting with •
Label in CAPS: then a specific, evidence-based action. Ground in published research.

STRICT RULES:
- Never use ** or * for emphasis. Use CAPS.
- No markdown formatting.
- No text outside the 4 sections."""

    try:
        response = model.generate_content(prompt)
        return response.text.strip()
    except Exception as e:
        log.warning("Gemini API call failed: %s", e)
        return None


def build_synopsis_slide(prs, rows, filters, fail_counts, header_text):
    synopsis_text = generate_gemini_synopsis(rows, filters, fail_counts)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    M = Inches(0.22)
    CONTENT_TOP = _draw_slide_chrome(slide, header_text + "  |  INSIGHTS")

    if not synopsis_text:
        add_text(slide, "[WARNING] Could not generate AI synopsis. Check your Gemini API key.",
                 _M_L, CONTENT_TOP, SLIDE_W - _M_L - M, Inches(1),
                 font="Poppins", size=13, color=C_RED_SCORE,
                 align=PP_ALIGN.CENTER, anchor="ctr")
        return

    synopsis_text = synopsis_text.replace("**", "").replace("*", "")

    section_keys = ["EXECUTIVE SUMMARY", "KEY STRENGTHS", "AREAS FOR IMPROVEMENT", "PRIORITY RECOMMENDATIONS"]
    sections = {}
    current_key, current_lines = None, []
    for line in synopsis_text.splitlines():
        stripped = line.strip()
        matched = False
        for key in section_keys:
            if stripped.upper().startswith(key):
                if current_key:
                    sections[current_key] = "\n".join(current_lines).strip()
                current_key, current_lines = key, []
                matched = True; break
        if not matched and current_key:
            current_lines.append(stripped)
    if current_key:
        sections[current_key] = "\n".join(current_lines).strip()

    COL_W = (SLIDE_W - _M_L - M - M) / 2
    COL_H = (SLIDE_H - CONTENT_TOP - M * 3) / 2
    positions = [
        (_M_L,             CONTENT_TOP),
        (_M_L + COL_W + M, CONTENT_TOP),
        (_M_L,             CONTENT_TOP + COL_H + M),
        (_M_L + COL_W + M, CONTENT_TOP + COL_H + M),
    ]
    section_colors = [C_PURPLE, C_TEAL, C_RED_SCORE, C_AMBER]

    for idx, key in enumerate(section_keys):
        body_text = sections.get(key, "No content generated.")
        cx, cy    = positions[idx]
        color     = section_colors[idx]

        add_rect(slide, cx, cy, COL_W, COL_H, fill=C_WHITE, line=C_HEADER_BDR, lw=0.5)
        add_rect(slide, cx, cy, COL_W, Inches(0.06), fill=color)
        add_text(slide, key, cx + Inches(0.2), cy + Inches(0.12),
                 COL_W - Inches(0.4), Inches(0.32),
                 font="Poppins", size=9, bold=True, color=color, anchor="t")

        body_clean = body_text.replace("**", "").replace("*", "")
        bullets    = [b.strip() for b in body_clean.split("•") if b.strip()]

        if not bullets:
            add_text(slide, body_clean, cx + Inches(0.2), cy + Inches(0.48),
                     COL_W - Inches(0.4), COL_H - Inches(0.6),
                     font="Poppins", size=8.5, color=C_BLACK, anchor="t", wrap=True)
        else:
            avail_h  = COL_H - Inches(0.6)
            bullet_h = avail_h / len(bullets)
            for bi, bullet in enumerate(bullets):
                by = cy + Inches(0.48) + bullet_h * bi
                label_part, rest = (bullet.split(":", 1) + [""])[:2] if ":" in bullet else ("", bullet)
                label_part, rest = label_part.strip(), rest.strip()

                txb = slide.shapes.add_textbox(cx + Inches(0.2), by,
                                               COL_W - Inches(0.4), bullet_h - Inches(0.05))
                tf = txb.text_frame
                tf.word_wrap = True
                from pptx.oxml.ns import qn as _qn
                bp = tf._txBody.find(_qn("a:bodyPr"))
                if bp is None:
                    from lxml import etree as _et
                    bp = _et.SubElement(tf._txBody, _qn("a:bodyPr"))
                bp.set("anchor", "t")
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT

                for text_content, is_label, is_dot in [
                    ("• ",                 False, True),
                    (label_part.upper() + ": " if label_part else None, True, False),
                    (rest,                 False, False),
                ]:
                    if text_content is None:
                        continue
                    run = p.add_run()
                    run.text           = text_content
                    run.font.name      = "Poppins"
                    run.font.size      = Pt(8.5)
                    run.font.bold      = is_label
                    run.font.color.rgb = color if (is_label or is_dot) else C_BLACK

    log.info("AI synopsis slide built.")


# ───────────────────────────────────────────────────────────────
# ANALYTICS SLIDE
# ───────────────────────────────────────────────────────────────

def build_analytics_slide(prs, rows, header_text, analysis_mode='average'):
    brand_map   = {}
    country_map = {}
    for r in rows:
        for field, fmap in [("brand", brand_map), ("country", country_map)]:
            v = r.get(field, "")
            if v:
                try:
                    fmap.setdefault(v, []).append(float(r["score"].replace("%", "")))
                except Exception:
                    pass

    brand_avgs   = {b: sum(v)/len(v) for b, v in brand_map.items()}
    country_avgs = {c: sum(v)/len(v) for c, v in country_map.items()}
    multi_brand   = len(brand_avgs) > 1
    multi_country = len(country_avgs) > 1

    if not multi_brand and not multi_country:
        log.info("Analytics slide skipped — single brand/country.")
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    M = Inches(0.22)
    CONTENT_TOP = _draw_slide_chrome(slide, header_text + "  |  ANALYTICS")
    bg_hex = "#FFFFFF"

    scores = []
    for r in rows:
        try:
            scores.append(float(r["score"].replace("%", "")))
        except Exception:
            pass

    charts_to_show = ["distribution"]
    if multi_brand:   charts_to_show.append("brand")
    if multi_country: charts_to_show.append("country")

    n_charts   = len(charts_to_show)
    chart_h    = SLIDE_H - CONTENT_TOP - M * 2
    chart_w    = (SLIDE_W - _M_L - M * (n_charts + 1)) / n_charts
    chart_h_in = chart_h / 914400
    chart_w_in = chart_w / 914400
    chart_paths = []

    for chart_type in charts_to_show:
        fig_path = DOWNLOAD_DIR / f"analytics_{chart_type}.png"
        fig, ax  = plt.subplots(figsize=(chart_w_in, chart_h_in), facecolor=bg_hex)
        ax.set_facecolor(bg_hex)
        for spine in ax.spines.values():
            spine.set_visible(False)
        ax.tick_params(colors="#444444", length=0)
        ax.yaxis.grid(True, color="#E8E8E8", linewidth=0.5, zorder=0)
        ax.set_axisbelow(True)

        if chart_type == "distribution":
            if len(scores) > 1 and np.var(scores) > 0:
                from scipy.stats import gaussian_kde
                kde    = gaussian_kde(scores)
                x_vals = np.linspace(0, 125, 500)
                y_vals = kde(x_vals) * len(scores) * 5
                ax.plot(x_vals, y_vals, color="#3C2A8C", linewidth=1.5, zorder=4)
                ax.fill_between(x_vals, 0, y_vals, where=(x_vals < 55),                     facecolor="#C00000", alpha=0.8, zorder=3)
                ax.fill_between(x_vals, 0, y_vals, where=((x_vals >= 55) & (x_vals < 85)),  facecolor="#E08C00", alpha=0.8, zorder=3)
                ax.fill_between(x_vals, 0, y_vals, where=(x_vals >= 85),                     facecolor="#009999", alpha=0.8, zorder=3)
                ymax = max(y_vals) if len(y_vals) > 0 else 1
            else:
                bins = list(range(0, 126, 5))
                counts, edges = np.histogram(scores, bins=bins)
                bar_cols = ["#C00000" if e < 55 else "#E08C00" if e < 85 else "#009999" for e in edges[:-1]]
                ax.bar(edges[:-1], counts, width=4.8, color=bar_cols, align="edge", zorder=3, edgecolor=bg_hex, linewidth=0.8)
                ymax = max(counts) if counts.any() else 1

            ax.axvspan(0,   55, alpha=0.04, color="#C00000", zorder=1)
            ax.axvspan(55,  85, alpha=0.04, color="#E08C00", zorder=1)
            ax.axvspan(85, 125, alpha=0.04, color="#009999", zorder=1)
            ax.set_ylim(0, ymax * 1.35)
            ax.text(27,  ymax*1.25, "MISSABLE",   color="#C00000", fontsize=7.5, ha="center", fontweight="bold")
            ax.text(70,  ymax*1.25, "WALLPAPER",  color="#E08C00", fontsize=7.5, ha="center", fontweight="bold")
            ax.text(103, ymax*1.25, "UNMISSABLE", color="#009999", fontsize=7.5, ha="center", fontweight="bold")
            ax.set_xlabel("Score %", color="#444444", fontsize=8)
            ax.set_ylabel("Approx. No. of Creatives", color="#444444", fontsize=8)
            ax.set_title("SCORE DISTRIBUTION", color="#3C2A8C", fontsize=9, fontweight="bold", pad=8)
            ax.set_xlim(0, 125)
            ax.xaxis.set_tick_params(labelcolor="#444444", labelsize=7)
            ax.yaxis.set_tick_params(labelcolor="#444444", labelsize=7)
            ax.yaxis.set_major_locator(plt.MaxNLocator(integer=True))

        elif chart_type in ("brand", "country"):
            raw_map = brand_map if chart_type == "brand" else country_map
            avg_map = brand_avgs if chart_type == "brand" else country_avgs
            if analysis_mode == "volume":
                sorted_data = sorted(raw_map.items(), key=lambda x: len(x[1]), reverse=True)[:15]
                labels   = [d[0].title() for d in sorted_data]
                values   = [len(d[1]) for d in sorted_data]
                bar_cols = ["#3C2A8C"] * len(values)
                x_label  = "Number of Creatives"
                title    = "BRAND VOLUME" if chart_type == "brand" else "COUNTRY VOLUME"
            else:
                sorted_data = sorted(avg_map.items(), key=lambda x: x[1], reverse=True)[:15]
                labels   = [d[0].title() for d in sorted_data]
                values   = [d[1] for d in sorted_data]
                bar_cols = ["#009999" if v >= 85 else "#E08C00" if v >= 55 else "#C00000" for v in values]
                x_label  = "Avg Score %"
                title    = "BRAND BENCHMARK" if chart_type == "brand" else "COUNTRY BENCHMARK"
            y_pos = np.arange(len(labels))
            if analysis_mode != "volume":
                ax.axvspan(0,   55, alpha=0.04, color="#C00000", zorder=1)
                ax.axvspan(55,  85, alpha=0.04, color="#E08C00", zorder=1)
                ax.axvspan(85, 125, alpha=0.04, color="#009999", zorder=1)
            bars = ax.barh(y_pos, values, color=bar_cols, zorder=3, height=0.6)
            for bar, val in zip(bars, values):
                ax.text(bar.get_width() + (max(values)*0.02), bar.get_y() + bar.get_height()/2,
                        f"{val:.0f}%" if analysis_mode != "volume" else str(val),
                        va="center", ha="left", color="#1A1A1A", fontsize=7, fontweight="bold")
            ax.set_yticks(y_pos)
            ax.set_yticklabels(labels, color="#444444", fontsize=7)
            ax.set_xlabel(x_label, color="#444444", fontsize=8)
            ax.set_title(title, color="#3C2A8C", fontsize=9, fontweight="bold", pad=8)
            ax.set_xlim(0, max(values)*1.2 if analysis_mode == "volume" else 125)
            ax.invert_yaxis()
            ax.xaxis.set_tick_params(labelcolor="#444444", labelsize=7)
            ax.xaxis.grid(True, color="#E8E8E8", linewidth=0.5, zorder=0)
            ax.yaxis.grid(False)

        plt.tight_layout(pad=0.5)
        plt.savefig(str(fig_path), dpi=150, bbox_inches="tight", facecolor=bg_hex, format="png")
        plt.close(fig)
        chart_paths.append(fig_path)

    for ci, fig_path in enumerate(chart_paths):
        cx = _M_L + ci * (chart_w + M)
        slide.shapes.add_picture(str(fig_path), cx, CONTENT_TOP, chart_w, chart_h)
        try:
            fig_path.unlink(missing_ok=True)
        except Exception:
            pass
    log.info("Analytics slide built (%d charts).", n_charts)
